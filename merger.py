import os
from pptx import Presentation
from slide_reader import get_slide_text_merger
from slide_writer import add_slide_with_text
from utils import print_file_grid

def main():
    output_dir = "./merged"
    os.makedirs(output_dir, exist_ok=True)

    output_filename = input("Enter your output file name (without .pptx): ").strip()
    if not output_filename:
        print("Output filename is required. Exiting.")
        return
    output_path = os.path.join(output_dir, f"{output_filename}.pptx")

    print_file_grid("./output")

    print("\nPaste the names of the input presentations, one per line.")
    print("Press Enter on an empty line when you're done:\n")

    input_prs_paths = []
    while True:
        path = input().strip()
        if not path:
            break
        if os.path.isfile(f"./output/{path}"):
            input_prs_paths.append(f"./output/{path}")
        else:
            print(f"File not found: {path}")

    if not input_prs_paths:
        print("No valid input files provided. Exiting.")
        return

    output_prs = Presentation()

    for prs_path in input_prs_paths:
        prs = Presentation(prs_path)

        for i in range(len(prs.slides)):
            slide = prs.slides[i]
            slide_text = get_slide_text_merger(slide)
            is_bold = i == 0 or i == len(prs.slides) - 1
            add_slide_with_text(output_prs, slide_text, bold=is_bold)

    output_prs.save(output_path)
    print(f"\nMerged presentation saved to: {output_path}")


main()