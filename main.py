from pptx import Presentation
import os, re

from slide_reader import get_slide_text, get_slide_type
from slide_writer import add_slide_with_text, create_new_presentation

def main():
    output_dir = "./output"
    os.makedirs(output_dir, exist_ok=True)

    source_presentation = Presentation("./source.pptx")
    output_file = None
    output_file_title = None

    for i in range(2, len(source_presentation.slides) - 1):
        slide = source_presentation.slides[i]
        should_save_prs = False
        slide_text = get_slide_text(slide)
        slide_type = get_slide_type(slide)

        if slide_type == "FIRST":
            output_file = create_new_presentation(source_presentation)
            output_file_title = re.sub(r"[^\w\d]+", "_", slide_text.strip())

        if slide_type == "LAST":
            should_save_prs = True

        add_slide_with_text(output_file, slide_text)

        if should_save_prs and output_file:
            filename = f"{output_dir}/{output_file_title}.pptx"
            output_file.save(filename)
            output_file = None
            output_file_title = None

if __name__ == "__main__":
    main()
