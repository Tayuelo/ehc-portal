import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def add_slide_with_text(prs, text):
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    textbox = new_slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.5),
        Inches(9),
        Inches(6.5),
    )

    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    for line in text.strip().split("\n"):
        if line.strip():
            p = text_frame.add_paragraph()
            p.text = line.strip()
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(65)

def extract_presentation(source: str, output_dir="./output"):
    prs = Presentation(source)
    slide_count = len(prs.slides)

    output_file = None
    prs_title = None
    song_started = False
    end_of_song_detected = False
    slide_text = ""
    first_song_slide_text = ""

    os.makedirs(output_dir, exist_ok=True)

    for i in range(2, slide_count - 1):
        slide = prs.slides[i]
        slide_text = ""
        end_of_song_detected = False

        full_slide_text = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    raw_text = run.text.strip()
                    if raw_text:
                        full_slide_text += raw_text + "\n"

                    if run.font.bold and not song_started:
                        first_song_slide_text = full_slide_text.strip()
                        prs_title = re.sub(r"[^\w\d]+", "_", first_song_slide_text)
                        prs_title = prs_title[:50]  # limit filename length
                        output_file = Presentation()
                        if output_file.slides:
                            output_file.slides.remove(output_file.slides[0])
                        song_started = True
                        print(f"🎵 Starting new song: {prs_title}")

                    if song_started and raw_text:
                        slide_text += raw_text + "\n"

                    if run.hyperlink and run.hyperlink.address and song_started:
                        end_of_song_detected = True

        if song_started and slide_text.strip():
            add_slide_with_text(output_file, slide_text)

        if end_of_song_detected and output_file:
            filename = f"{output_dir}/{prs_title}.pptx"
            output_file.save(filename)
            print(f"✅ Saved: {filename}")
            output_file = None
            prs_title = None
            song_started = False
            slide_text = ""
            first_song_slide_text = ""

extract_presentation("./source.pptx")