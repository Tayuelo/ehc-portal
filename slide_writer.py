from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx import Presentation

def create_new_presentation(source_presentation):
    new_prs = Presentation()
    new_prs.slide_width = source_presentation.slide_width
    new_prs.slide_height = source_presentation.slide_height
    return new_prs

def add_slide_with_text(prs, text, horizontal_padding=Inches(1)):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    width = prs.slide_width - 2 * horizontal_padding
    height = prs.slide_height
    left = horizontal_padding
    top = 0

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER

    run = paragraph.add_run()
    run.text = text.strip()
    run.font.size = Pt(65)
